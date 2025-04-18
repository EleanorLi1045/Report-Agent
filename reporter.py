
import os
import logging
from ring_classifier import get_ring_progress_from_status
from datetime import datetime, timezone
import pandas as pd
import requests
from requests.auth import HTTPBasicAuth
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
import argparse
from datetime import datetime
import logging
import sys
import argparse

# ADO Configuration
ORGANIZATION    = "domoreexp"
PROJECT         = "MSTeams"

#QUERY_ID        = "06e3344a-b5b8-40b4-a4e3-9a8a72f5b034"
#QUERY_ID        = "1cc233d6-00b3-4ca2-a846-b15ea0a9e0aa"
#QUERY_ID        = "b9458d82-7f05-44e4-84e6-71ca92b54089"
#QUERY_ID        = "44774f16-3f94-4bb7-9724-0e7eef8e5de1"
#QUERY_ID        = "ca0c709b-bcfd-423a-b3b5-687a7acc1839" # GA item tree relation
#ORGANIZATION    = "onedrive"
#PROJECT         = "Path"
#QUERY_ID        = "f2a5ca93-4594-4bbd-900f-633f67ecb96b"
ai_enabled       = False

REQUIRED_FIELD = {
        "System.Title",
        "Microsoft.VSTS.Scheduling.TargetDate",
        "MicrosoftTeamsCMMI.Ring1TargetDate",
        "MicrosoftTeamsCMMI.Ring2TargetDate",
        "MicrosoftTeamsCMMI.Ring3TargetDate",
        "MicrosoftTeamsCMMI.Ring4TargetDate",
        "MicrosoftTeamsCMMI.StatusTweet",
        "MicrosoftTeamsCMMI.PMPlanning",
        "MicrosoftTeamsCMMI.DesignPlanning",
        "MicrosoftTeamsCMMI.DevPlanning",
        "System.Tags",
        "MicrosoftTeamsCMMI.ShiproomStatusColor",
        "System.State"        
}
COLUMN_CONFIG = {
    2: ("System.Title", 8),
    3: ("EM", 7),
    4: ("PM", 7),
    10: ("Microsoft.VSTS.Scheduling.TargetDate", 6),
    11: ("Microsoft.VSTS.Scheduling.TargetDate", 6),
    12: ("MicrosoftTeamsCMMI.Ring1TargetDate", 6),
    13: ("MicrosoftTeamsCMMI.Ring2TargetDate", 6),
    14: ("MicrosoftTeamsCMMI.Ring3TargetDate", 6),
    15: ("MicrosoftTeamsCMMI.Ring3TargetDate", 6),
    16: ("MicrosoftTeamsCMMI.Ring3TargetDate", 6),
    17: ("MicrosoftTeamsCMMI.Ring3TargetDate", 6),
    18: ("MicrosoftTeamsCMMI.Ring4TargetDate", 6),
    19: ("MicrosoftTeamsCMMI.StatusTweet", 7),
}
DATE_FIELDS = {
    "Microsoft.VSTS.Scheduling.TargetDate",
    "MicrosoftTeamsCMMI.Ring1TargetDate",
    "MicrosoftTeamsCMMI.Ring2TargetDate",
    "MicrosoftTeamsCMMI.Ring3TargetDate",
    "MicrosoftTeamsCMMI.Ring4TargetDate"
}

def set_date_cell_with_highlight(table, row_index, col_index, fields, field_name, font_size):
    try:
        cell = table.cell(row_index, col_index)
        raw_value = fields.get(field_name, "")
        formatted_value = raw_value

        if raw_value:
            try:
                parsed_date = pd.to_datetime(raw_value)
                if parsed_date.date() < datetime.now(timezone.utc).date():
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = MUTED_YELLOW
                formatted_value = parsed_date.strftime('%m/%d')
            except Exception as e:
                print(f"⚠️ Date error for {field_name}: {raw_value} → {e}")

        cell.text = formatted_value

        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.name = "Helvetica Neue"
                run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

    except Exception as e:
        print(f"⚠️ Error processing date field {field_name}: {e}")

def read_pat(file_path="PAT.txt"):
    with open(file_path, "r") as f:
        return f.read().strip()

def fetch_work_item_details(pat, ids):
    print(f"🔍 Fetching details for {len(ids)} work items...")
    all_details = []
    batch_size = 100
    for i in range(0, len(ids), batch_size):
        id_batch = ids[i:i+batch_size]
        ids_str = ",".join(str(i) for i in id_batch)

        url = f"https://{ORGANIZATION}.visualstudio.com/_apis/wit/workitems?ids={ids_str}&api-version=6.0"
        response = requests.get(url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
        response.raise_for_status()
        
        print(f"✅ Fetching ADI Items batch {i+1}-{i+len(id_batch)}")

        all_details.extend(response.json().get("value", []))

    return all_details

def fetch_work_items_flat(pat,QUERY_ID):
    print("📡 Fetching work item with flat structures...")
    url = f"https://{ORGANIZATION}.visualstudio.com/{PROJECT}/_apis/wit/wiql/{QUERY_ID}?api-version=6.0"
    response = requests.get(url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
    response.raise_for_status()
    work_item_refs = response.json().get("workItems", [])
    ids = [item["id"] for item in work_item_refs]
    print(f"✅ Fetched {len(ids)} work item references. {ids}")
    return fetch_work_item_details(pat, ids)


def fetch_work_items_tree(pat,QUERY_ID):
    print("📡 Fetching work items with tree structure...")
    url = f"https://{ORGANIZATION}.visualstudio.com/{PROJECT}/_apis/wit/wiql/{QUERY_ID}?api-version=6.0"
    try:
        response = requests.get(url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
        response.raise_for_status()
        wiql_result = response.json()
        feature_ids = []
        work_item_relations = wiql_result.get("workItemRelations", [])

        for relation in work_item_relations:
            target_id = relation.get("target", {}).get("id")
            if target_id:
                try:
                    item_details_url = f"https://{ORGANIZATION}.visualstudio.com/{PROJECT}/_apis/wit/workitems/{target_id}?api-version=6.0"
                    details_response = requests.get(item_details_url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
                    details_response.raise_for_status()
                    item_details = details_response.json()

                    item_type = item_details.get("fields", {}).get("System.WorkItemType")
                    if item_type == "Feature":
                        feature_ids.append({"id": target_id})  # Append a dictionary instead of just the ID
                        print(f"✅ Found Feature ID: {target_id}")

                except requests.exceptions.RequestException as e:
                    print(f"⚠️ Error fetching details for work item {target_id}: {e}")
                except Exception as e:
                    print(f"⚠️ Error processing details for work item {target_id}: {e}")

        print(f"✅ Found {len(feature_ids)} feature IDs.")
        return feature_ids

    except requests.exceptions.RequestException as e:
        print(f"❌ Fatal error fetching work items: {e}")
        return []


def fetch_item_details_with_parent(work_items, pat, fetch_parent):
    print("📥 Fetching work item details...")

    rows = []
    for item in work_items:
        try:
            work_item_id = item["id"]
            url = f"https://{ORGANIZATION}.visualstudio.com/{PROJECT}/_apis/wit/workitems/{work_item_id}?$expand=relations&api-version=6.0"
            res = requests.get(url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
            res.raise_for_status()
            data = res.json()
            fields = data.get("fields", {})
            row = {field: fields.get(field, "") for field in REQUIRED_FIELD}
            row["Work Item Number"] = work_item_id
            title = fields.get("System.Title", "")
            print(f"✅ Processing work item {work_item_id}-{title} with fetch parent: {fetch_parent}")
            #print (fields)
            # Check for parent link
            parent_id    = None
            if fetch_parent:
                for relation in data.get("relations", []):
                    if relation.get("rel") == "System.LinkTypes.Hierarchy-Reverse":
                        parent_url = relation.get("url", "")
                        parent_id = parent_url.split("/")[-1]
                        break

                if parent_id:
                    # Fetch parent title
                    parent_url = f"https://{ORGANIZATION}.visualstudio.com/{PROJECT}/_apis/wit/workitems/{parent_id}?api-version=6.0"
                    parent_res = requests.get(parent_url, headers={"Content-Type": "application/json"}, auth=HTTPBasicAuth('', pat))
                    parent_res.raise_for_status()
                    parent_data = parent_res.json()
                    parent_title = parent_data.get("fields", {}).get("System.Title", "Unknown")
                    row["Parent ID"]        = parent_id
                    row["Parent Title"]     = parent_title
                    row["PM"]               = fields.get('MicrosoftTeamsCMMI.PMOwner').get('displayName', 'Unknown')
                    row["EM"]               = fields.get('MicrosoftTeamsCMMI.EMOwner').get('displayName', 'Unknown')
                    print(f"🔗 Work item {work_item_id}-{title} → Parent: {parent_title} ({parent_id})")
                else:
                    row["Parent ID"]        = ""
                    row["Parent Title"]     = ""
                    row["PM"]               = fields.get('MicrosoftTeamsCMMI.PMOwner').get('displayName', 'Unknown')
                    row["EM"]               = fields.get('MicrosoftTeamsCMMI.EMOwner').get('displayName', 'Unknown')
                    print(f"{fields.get('MicrosoftTeamsCMMI.PMOwner').get('displayName', 'Unknown')} - ⚠️ Work item {work_item_id}-{title} has no parent.")
                
            rows.append(row)

        except Exception as e:
            print(f"⚠️ Failed work item {item['id']}: {e}")
    return rows

def group_items_by_parent(work_items):
    print("\n📊 Grouping work items by parent epic...")

    grouped = {}

    for item in work_items:
        parent_title = item.get("Parent Title", "No Parent")
        parent_id = item.get("Parent ID", "")
        key = f"{parent_title} ({parent_id})" if parent_id else "No Parent"

        if key not in grouped:
            grouped[key] = []
        grouped[key].append(item)

    for parent, children in grouped.items():
        print(f"\n🧩 Parent Epic: {parent}")
        for child in children:
            print(f"   - [{child['Work Item Number']}] {child.get('System.Title', 'No Title')}")

    return grouped

def categorize_by_workstream(items):
    print("🔄 Categorizing work items by workstream...")
    workstreams = ["Workstream 1", "Workstream 2", "Workstream 3", "Workstream 4","output_driven_meetings", "Prism and AI Quality"]
    
    workstream_data = {ws: [] for ws in workstreams}
    workstream_data["Other Workstreams"] = []
    workstream_data["all"]               = []

    for item in items:
        tags_raw    = item.get("System.Tags", "")
        tags        = [t.strip() for t in tags_raw.split(";") if t.strip()]
        matched     = False
        for ws in workstreams:
            if ws in tags:
                workstream_data[ws].append(item)
                workstream_data["all"].append(item)
                matched = True
                break
        if not matched:
            workstream_data["all"].append(item)
            workstream_data["Other Workstreams"].append(item)

    for ws, data in workstream_data.items():
        print(f"📊 {ws}: {len(data)} items")
    return workstream_data

def hex_to_rgb_color(hex_color):
    """
    Converts a hex color string like '#B6D7A8' to an RGBColor object.
    """
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16))

def set_planning_status_emoji(table, row_index, column_index, fields, field_key, start_row=0):
    """
    Sets an emoji and background color in a table cell based on planning status.
    """
    cell = table.cell(row_index + start_row, column_index)
    planning_status = fields.get(field_key, "").strip()

    green_statuses = [
        "1-Pager Ready for Review", "Signed Off", "Full Spec Ready for Review",
        "Complete", "N/A", "Spec: Signed off", "Completed"
    ]
    yellow_statuses = [
        "Drafting", "In Progress", "Spec: In review", "Spec: Draft", "In Review"
    ]

    # Determine emoji and fill color
    if planning_status in green_statuses:
        emoji = "🟢"
        fill_color = "#B6D7A8"  # Soft Green
    elif planning_status in yellow_statuses:
        emoji = "🟡"
        fill_color = "#D8B56A"  # Muted Yellow
    elif planning_status == "Required":
        emoji = "🔴"
        fill_color = "#D98880"  # Muted Red
    elif planning_status == "":
        emoji = "🟠"
        fill_color = "#D3D3D3"  # Light Gray (Not Started)
    else:
        emoji = "⚠️"
        fill_color = "#FFFFFF"  # White background for unexpected values
        print(f"⚠️ Unexpected {field_key} value: '{planning_status}'")

    # Lets NOT Set emoji text as we are using color code
    #cell.text = emoji 

    # Center align and shrink font size
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(7)

    # Apply background fill color
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb_color(fill_color)

MUTED_YELLOW    = RGBColor(0xD8, 0xB5, 0x6A)
SOFT_GREEN      = RGBColor(0xB6, 0xD7, 0xA8)
WHITE           = RGBColor(255, 255, 255)
MUTED_RED        = RGBColor(0xD9, 0x88, 0x80)



def getShipRoomStatusColor(fields):
    status = fields.get("MicrosoftTeamsCMMI.ShiproomStatusColor", "").lower()
    if status == "green":
        return RGBColor(0, 176, 80)  # Green
    elif status == "red":
        return RGBColor(255, 0, 0)  # Red
    elif status == "orange":
        return RGBColor(255, 165, 0)  # Orange
    else:
        return RGBColor(128, 128, 128)  # Gray

def fill_template_ppt(workstream_data, template_path, output_path):
    print(f"🎨 Filling PowerPoint: {template_path} -> {output_path}")
    prs = Presentation(template_path)
    sorted_config = sorted(COLUMN_CONFIG.items())
    column_numbers = [int(c) for c in COLUMN_CONFIG.keys()]
    field_names = [v[0] for _, v in sorted_config]
    font_sizes = [v[1] for _, v in sorted_config]
    for slide in prs.slides:
        title_shape = next((s for s in slide.shapes if s.has_text_frame), None)
        if not title_shape:
            continue
        title = title_shape.text.strip()
        if title not in workstream_data:
            continue

        df = pd.DataFrame(workstream_data[title])
        table = next((shape.table for shape in slide.shapes if shape.has_table), None)
        if not table:
            continue
       
        start_row = 2
        for i, (_, row_data) in enumerate(df.iterrows()):
            if i + start_row >= len(table.rows):
                break

            fields = {key: str(row_data.get(key, "")).strip() for key in df.columns}
            ado_title = fields.get("System.Title", "")

            set_planning_status_emoji(table, i, 5, fields, "MicrosoftTeamsCMMI.PMPlanning", start_row) #not generic , please fix later
            set_planning_status_emoji(table, i, 6, fields, "MicrosoftTeamsCMMI.DesignPlanning", start_row) #not generic , please fix later
            set_planning_status_emoji(table, i, 7, fields, "MicrosoftTeamsCMMI.DevPlanning", start_row) #not generic , please fix later

            state = fields.get("System.State", "")
            if state != "Proposed":
                cell = table.cell(i + start_row, 8) #not generic , please fix later
                cell.fill.solid()
                cell.fill.fore_color.rgb = SOFT_GREEN
            
            fill_parent = True #not generic , please fix later
            if fill_parent:
                parent    = fields.get("Parent Title", "Parent Title")
                cell      = table.cell(i + start_row, 0) #not generic , please fix later
                cell.text = parent #not generic , please fix later
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)
                        run.font.name = "Helvetica Neue"
                        run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
            
            cell = table.cell(i + start_row, 4) #not generic , please fix later
            cell.fill.solid()
            cell.fill.fore_color.rgb = getShipRoomStatusColor(fields)
            

            # Ring3 target date
            #r3_date = fields.get("System.Ring3TargetDate", "")
            #cell = table.cell(i + start_row, 13) #not generic , please fix later
            #cell.text = r3_date
            #for paragraph in cell.text_frame.paragraphs:
            #    paragraph.alignment = PP_ALIGN.CENTER
            #    for run in paragraph.runs:
            #        run.font.size = Pt(7)

            #find in wich RIng this features has released
            status_tweet_for_ai = fields.get("MicrosoftTeamsCMMI.StatusTweet", "")
            ring_progress = ""
            if ai_enabled :
                ring_progress  = get_ring_progress_from_status(status_tweet_for_ai)
                print(f"👨‍💻 {fields.get("MicrosoftTeamsCMMI.StatusTweet", "")} - Release Status: {ring_progress}")
           

            # Process all configured fields
            for logical_idx, field_name in enumerate(field_names):
                col_idx = column_numbers[logical_idx] - 1
                font_size = font_sizes[logical_idx]

                if field_name in DATE_FIELDS:
                    row_index = i + start_row
                    col_index = col_idx
                      
                    try:
                        cell = table.cell(row_index, col_index)
                        raw_value = fields.get(field_name, "")
                        formatted_value = raw_value

                        if raw_value:
                            try:
                                parsed_date = pd.to_datetime(raw_value, utc=True)
                                if (parsed_date.date() < datetime.now(timezone.utc).date()) and ai_enabled:
                                    # Release Date is in Past
                                    # ring_progress
                                    any_date_red = False
                                    if field_name == "Microsoft.VSTS.Scheduling.TargetDate":
                                        if ring_progress["ring_0"] == 100:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = SOFT_GREEN
                                        else:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = MUTED_RED
                                            any_date_red = True
                                    if field_name == "MicrosoftTeamsCMMI.Ring1TargetDate":
                                        if ring_progress["ring_1"] == 100:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = SOFT_GREEN
                                        else:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = MUTED_RED
                                            any_date_red = True
                                    if field_name == "MicrosoftTeamsCMMI.Ring2TargetDate":
                                        if ring_progress["ring_2"] == 100:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = SOFT_GREEN
                                        else:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = MUTED_RED
                                            any_date_red = True
                                    if field_name == "MicrosoftTeamsCMMI.Ring3TargetDate":
                                        if ring_progress["ring_3"] == 100:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = SOFT_GREEN
                                        else:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = MUTED_RED
                                            any_date_red = True
                                    if field_name == "MicrosoftTeamsCMMI.Ring4TargetDate":
                                        if ring_progress["ring_4"] == 100:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = SOFT_GREEN
                                        else:
                                            cell.fill.solid()
                                            cell.fill.fore_color.rgb = MUTED_RED
                                            any_date_red = True
                                    if any_date_red:
                                        shipRoomStatusCell = table.cell(i + start_row, 4) #not generic , please fix later
                                        shipRoomStatusCell.fill.solid()
                                        shipRoomStatusCell.fill.fore_color.rgb = RGBColor(255, 0, 0)

                                formatted_value = parsed_date.strftime('%m/%d')
                            except Exception as e:
                                print(f"⚠️ Date error for {ado_title} - {field_name}: {raw_value} → {e}")

                        cell.text = formatted_value

                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(font_size)
                                run.font.name = "Helvetica Neue"
                                run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)

                    except Exception as e:
                        print(f"⚠️ Error processing {ado_title} date field {field_name}: {e}")
                else:
                    try:
                        cell = table.cell(i + start_row, col_idx)
                        value = fields.get(field_name, "")
                        cell.text = value

                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(font_size)
                                run.font.name = "Helvetica Neue"
                                run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
                    except Exception as e:
                        print(f"⚠️ PPT fill error (row {i}, col {col_idx}): {e}")
             
    prs.save(output_path)
    print("✅ PowerPoint saved!")

def rgbcolor_to_argb_hex(rgb_color: RGBColor) -> str:
    """Converts pptx RGBColor to ARGB hex format for Excel."""
    return "FF{:02X}{:02X}{:02X}".format(rgb_color[0], rgb_color[1], rgb_color[2])

def process_ppt_to_excel_with_color(ppt_path: str, excel_template_path: str, output_excel_path: str) -> None:
    logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")

    print(f"🔄 Starting the PowerPoint to Excel process {ppt_path}.")

    if not os.path.exists(ppt_path):
        print(f"⚠️ PPTX file not found at: {ppt_path}")
        return

    if not os.path.exists(excel_template_path):
        print(f"⚠️ Excel template not found at: {excel_template_path}")
        return

    prs = Presentation(ppt_path)
    wb = load_workbook(excel_template_path)

    workstream_titles = {
        "Workstream 1": "Workstream 1",
        "Workstream 2": "Workstream 2",
        "Workstream 3": "Workstream 3",
        "Workstream 4": "Workstream 4",
        "output_driven_meetings": "output_driven_meetings",
        "Prism and AI Quality": "Prism and AI Quality",
        "Other Workstreams": "Other Workstreams",
        "all": "all"
    }


    for slide in prs.slides:
        if not slide.shapes.title:
            continue
        title = slide.shapes.title.text.strip()
        sheet_name = workstream_titles.get(title)

        if not sheet_name or sheet_name not in wb.sheetnames:
            continue

        table_shape = next((shape for shape in slide.shapes if shape.has_table), None)
        if not table_shape:
            print(f"⚠️ No table found in slide: '{title}'. Skipping...")
            continue

        table = table_shape.table
        sheet = wb[sheet_name]

        for row_idx in range(2, len(table.rows)):  # Start from row 3 in Excel
            for col_idx, ppt_cell in enumerate(table.rows[row_idx].cells):
                value = ppt_cell.text.strip()
                excel_cell = sheet.cell(row=row_idx + 1, column=col_idx + 1)
                excel_cell.value = value

                # Copy background color if available
                try:
                    ppt_fill = ppt_cell.fill
                    if ppt_fill.type == 1 and ppt_fill.fore_color.type == 1:  # solid + RGB
                        rgb = ppt_fill.fore_color.rgb
                        if rgb:
                            hex_color = rgbcolor_to_argb_hex(rgb)
                            excel_cell.fill = PatternFill(
                                start_color= hex_color,
                                end_color  = hex_color,
                                fill_type  = "solid"
                            )
                except Exception as e:
                    print(f"⚠️ Error copying cell color at row {row_idx+1}, col {col_idx+1}: {e}")

        print(f"📊 Successfully added data (with color) to sheet: '{sheet_name}'")

    wb.save(output_excel_path)
    print(f"✅ Excel file successfully saved to: {output_excel_path}")

def create_ppt_from_grouped_items(grouped, required_fields, output_path="grouped_epics.pptx"):
    print("\n📽️ Creating PPT with epics and their items...")
    prs = Presentation()

    for parent, items in grouped.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = f"Epic: {parent}"

        rows = len(items) + 1
        cols = len(required_fields)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Set header row
        for col_index, field in enumerate(required_fields):
            table.cell(0, col_index).text = field.split('.')[-1]

        # Fill data
        for row_index, item in enumerate(items, start=1):
            for col_index, field in enumerate(required_fields):
                value = str(item.get(field, ""))
                table.cell(row_index, col_index).text = value

    prs.save(output_path)
    print(f"✅ PPT saved to: {output_path}")




def merge_consecutive_cells_in_first_column(excel_filepath):
    """
    Merges consecutive cells with the same value in the first column of an Excel file.

    Args:
        excel_filepath (str): The path to the Excel file.
    """
    try:
        workbook = load_workbook(excel_filepath)
        sheet = workbook.active  # Assuming you want to work with the active sheet

        if sheet.max_row < 3:
            print("Excel file has less than 3 rows. No merging will be performed.")
            return

        merge_ranges = []
        start_row = 3
        current_value = sheet.cell(row=start_row, column=1).value

        for row_num in range(4, sheet.max_row + 1):
            cell_value = sheet.cell(row=row_num, column=1).value
            if cell_value == current_value:
                continue
            else:
                if row_num - start_row > 1:
                    merge_ranges.append(f'A{start_row}:A{row_num - 1}')
                start_row = row_num
                current_value = cell_value

        # Check for the last sequence
        if sheet.max_row - start_row + 1 > 1:
            merge_ranges.append(f'A{start_row}:A{sheet.max_row}')

        for merge_range in merge_ranges:
            sheet.merge_cells(merge_range)

        workbook.save(excel_filepath)
        print(f"Consecutive cells merged in the first column of '{excel_filepath}'.")

    except FileNotFoundError:
        print(f"Error: File not found at '{excel_filepath}'.")
    except Exception as e:
        print(f"An error occurred: {e}")

def merge_consecutive_cells_in_first_column_pptx(pptx_filepath):
    """
    Updates every slide of a PPTX presentation where a table is present.
    For each table found, it merges consecutive cells with the same value
    in the first column (starting from the third row) and removes
    duplicate content, keeping only the value from the first non-empty cell
    of the merged range. It avoids merging ranges that overlap with existing merged cells.
    Slides without tables are skipped.

    Args:
        pptx_filepath (str): The path to the PPTX file.
    """
    try:
        prs = Presentation(pptx_filepath)
        for slide in prs.slides:
            table_found = False
            for shape in slide.shapes:
                if shape.has_table:
                    table_found = True
                    table = shape.table
                    rows = table.rows
                    cols = table.columns
                    if len(rows) < 3:
                        print(f"Table in slide index {prs.slides.index(slide) + 1} has less than 3 rows. Skipping.")
                        continue

                    merge_origins = set()
                    for r_idx, row in enumerate(rows):
                        first_cell = row.cells[0]
                        if first_cell.is_merge_origin:
                            merge_origins.add(r_idx)

                    i = 2  # Start from the 3rd row (index 2)
                    while i < len(rows):
                        start_row_idx = i
                        start_cell = rows[start_row_idx].cells[0]

                        # Skip if the current row contains a merge origin (part of an existing merge)
                        is_part_of_existing_merge = False
                        for origin_row in merge_origins:
                            # We need a way to determine the span of the merge from the origin.
                            # Without row_span, this is tricky. Let's try a simpler approach:
                            # If the current row is a merge origin, skip it.
                            if start_row_idx == origin_row:
                                is_part_of_existing_merge = True
                                break

                        if is_part_of_existing_merge:
                            i += 1
                            continue

                        if start_cell is not None and start_cell.text is not None:
                            current_value = start_cell.text
                            end_row_idx = start_row_idx
                            for j in range(start_row_idx + 1, len(rows)):
                                next_cell = rows[j].cells[0]
                                is_next_part_of_existing_merge = False
                                for origin_row in merge_origins:
                                    if j == origin_row:
                                        is_next_part_of_existing_merge = True
                                        break

                                if is_next_part_of_existing_merge:
                                    break
                                cell_value = next_cell.text if next_cell is not None else None
                                if cell_value == current_value:
                                    end_row_idx = j
                                else:
                                    break

                            if end_row_idx > start_row_idx:
                                first_cell = rows[start_row_idx].cells[0]
                                if first_cell is not None:
                                    first_cell_value = first_cell.text
                                    try:
                                        
                                        table.cell(start_row_idx, 0).merge(table.cell(end_row_idx, 0))
                                        
                                        cell = table.cell(start_row_idx, 0)
                                        cell.text = first_cell_value
                                        cell.fill.solid()
                                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = Pt(8)
                                                run.font.name = "Helvetica Neue"
                                                run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
                                    except Exception as e:
                                        print(f"Error merging cells in row {start_row_idx + 1} to {end_row_idx + 1}: {e}")
                                i = end_row_idx + 1
                            else:
                                i += 1
                        else:
                            i += 1
            if table_found:
                print(f"Processed table in slide index {prs.slides.index(slide) + 1}: Merged cells in the first column and removed duplicates (handling existing merges).")
            else:
                print(f"No table found in slide index {prs.slides.index(slide) + 1}. Skipping.")

        prs.save(pptx_filepath)
        print(f"Consecutive cells merged and duplicate content removed in the first column of tables in '{pptx_filepath}' (existing merges handled).")

    except FileNotFoundError:
        print(f"Error: File not found at '{pptx_filepath}'.")
    except Exception as e:
        print(f"An error occurred: {e}")


logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def main():
    try:
        timestamp       = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_file        = f"CMD AI Workstream Updates_{timestamp}.pptx"
        excel_file      = f"CMD AI Workstream Updates_{timestamp}.xlsx"
        pat             = read_pat()

        # --- Combined argument parsing ---
        parser = argparse.ArgumentParser(description="Run the reporter with optional AI and tree structure options, and a mandatory Query ID.")
        parser.add_argument(
            "--ai",
            type=str,
            default="false",
            choices=["true", "false"],
            help="Set to true to enable AI functionality.  Default is false."
        )
        parser.add_argument(
            "--isTree",
             type=str,
             default="false",
             choices=["true", "false"],
             help="Fetch work items as a tree structure (true) or flat list (false). Default is false."
         )
        parser.add_argument(
             "--query_id",
             type=str,
             required=True,
             help="The ADO Query Id of the Azure DevOps query to fetch work items from."
         )
        args = parser.parse_args()

        query_id   = args.query_id
        ai_enabled = args.ai.lower() == "true"
        is_tree    = args.isTree.lower() == "true"

        print(f"Work items are fetched as a {'tree' if is_tree else 'flat list'}.")
             
        if ai_enabled:
             print(f"AI based shiproom color is {'enabled' if ai_enabled else 'disabled'}.")
        
        if query_id:
             logging.info(f"Using Query ID from command line: {query_id}")
        else:
             print("Query ID was not provided via command line.")
             query_id_input = input("Please enter the Query ID: ")
             if query_id_input:
                 query_id = query_id_input.strip()
                 logging.info(f"Using Query ID provided by user: {query_id}")
             else:
                 print("Error: Query ID cannot be empty. Exiting.")
                 sys.exit(1)

        if is_tree:
             print(f"Fetching work items as a tree structure using the query ID {query_id}.")
             work_items = fetch_work_items_tree(pat,query_id) # Assuming fetch_work_items_tree can accept query_id
        else:
             print(f"Fetching work items as a flat structure using the query ID {query_id}.")
             work_items = fetch_work_items_flat(pat,query_id)   # Assuming fetch_work_items_flat can accept query_id

        
         
         # --- End of argument parsing ---
        
        details         = fetch_item_details_with_parent(work_items, pat, True)
        workstream_data = categorize_by_workstream(details)
        #grouped        = group_items_by_parent(details)

        fill_template_ppt(workstream_data, "templates/facilitator_ga.pptx", ppt_file)

        process_ppt_to_excel_with_color(ppt_path=ppt_file, excel_template_path="templates/facilitator_ga.xlsx", output_excel_path=excel_file)
        
        #lets merge JTBD
        merge_consecutive_cells_in_first_column(excel_file)
        merge_consecutive_cells_in_first_column_pptx(ppt_file)
        
         #create_ppt_from_grouped_items(grouped, REQUIRED_FIELD, output_path="grouped_epics.pptx")

    except Exception as e:
         logging.error(f"❌ Fatal error in main function: {e}")
         print("Usage: python your_script_name.py --query_id <your_query_id> [--ai <true|false>] [--isTree <true|false>]")
         print("Please provide the Query ID.")
         sys.exit(1)

if __name__ == "__main__":
    main()

